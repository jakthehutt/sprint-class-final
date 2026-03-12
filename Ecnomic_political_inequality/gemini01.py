import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_str):
    """Convert hex color string (e.g. '#DB2777' or 'DB2777') to pptx RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def add_axis_titles(chart, x_title=None, y_title=None):
    """Add X and Y axis titles to a given chart if supported."""
    if hasattr(chart, 'category_axis') and x_title:
        try:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = x_title
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except Exception:
            pass
    if hasattr(chart, 'value_axis') and y_title:
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_title
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except Exception:
            pass

def create_inequality_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Standard layouts
    title_slide_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]

    # Color Palette from HTML
    PINK = '#DB2777'
    BLUE = '#3B82F6'
    EMERALD = '#10B981'
    PURPLE = '#8B5CF6'
    ORANGE = '#F59E0B'
    RED = '#EF4444'

    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "The Architecture of Inequality in Austria"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Structural Divides, Economic Realities, and Policy Solutions\n\nCourse: Economic and Political Inequality\nContext: Vienna, Austria"
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(20)

    # -------------------------------------------------------------------------
    # SLIDE 2: Inheritance & The Squeezed Middle
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "1. The 'Inheritance Economy' & The Squeezed Middle"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(PINK)
    
    p = tf.add_paragraph()
    p.text = "The traditional pathway to middle-class stability via labor is eroding. Maintaining class status now relies heavily on intergenerational wealth transfers. Based on the HFCS 2023, nearly 41% of Austrian households received an inheritance, and almost half of all total wealth in Austria is inherited. Migrants and non-inheritors are systematically excluded."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(PINK)
    
    p = tf.add_paragraph()
    p.text = "• Introduce a progressive inheritance tax (exemptions up to €1M, affecting 0.2-1% of heirs, generating up to €1.8B for the state).\n• Banning tax-exempt transfer of corporate equity.\n• Fund a universal 'starting capital' endowment for all young adults."
    p.font.size = Pt(14)

    # Chart: Doughnut
    chart_data = CategoryChartData()
    chart_data.categories = ['Top 10% Inheritors', 'Remaining Native Pop.', 'Migrant Background Pop.']
    chart_data.add_series('Share of Total Inherited Wealth', (65, 30, 5))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.plots[0].has_data_labels = True
    
    # Try coloring series 0 points
    colors = [PINK, BLUE, EMERALD]
    try:
        for idx, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(colors[idx])
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 3: Wealth vs Income Paradox
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "2. The Wealth vs. Income Paradox"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(EMERALD)
    
    p = tf.add_paragraph()
    p.text = "Austria presents a distinct paradox: its robust welfare state and progressive income tax keep income inequality remarkably low (Gini 0.28). However, the absence of capital and wealth taxes makes its net wealth inequality among the highest in the developed world (Gini 0.76). According to the OeNB, the Top 10% command >50% of the wealth."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(EMERALD)
    
    p = tf.add_paragraph()
    p.text = "• Implement a 0.5% - 1.5% net wealth tax on assets > €1M.\n• Equalize capital gains tax with top marginal income tax brackets to halt disproportionate capital accumulation."
    p.font.size = Pt(14)

    # Chart: Bar
    chart_data = CategoryChartData()
    chart_data.categories = ['Income Gini (After Taxes)', 'Wealth Gini (Gross Assets)']
    chart_data.add_series('Gini Coefficient', (0.28, 0.76))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.plots[0].has_data_labels = True
    add_axis_titles(chart, x_title='Gini Indicator', y_title='Coefficient (0-1)')
    
    colors = [EMERALD, PINK]
    try:
        for idx, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(colors[idx])
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 4: Part-Time Trap & Gender Pay Gap
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "3. The Part-Time Trap & Gender Pay Gap"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(ORANGE)
    
    p = tf.add_paragraph()
    p.text = "Austria maintains an 18% unadjusted gender pay gap (Statistics Austria, 2024). Deeply rooted structural norms force women into unpaid care work. Consequently, female part-time employment sits at an extreme 51.1% (vs 17% for men), cascading into a devastating 40.7% gender pension gap later in life."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(ORANGE)
    
    p = tf.add_paragraph()
    p.text = "• Fully fund and systematically expand right-to-childcare infrastructure from age one nationwide.\n• Introduce non-transferable parental leave quotas ('daddy months') to forcefully shift cultural caregiving norms."
    p.font.size = Pt(14)

    # Chart: Column
    chart_data = CategoryChartData()
    chart_data.categories = ['Men Full-Time', 'Women Full-Time', 'Women Part-Time']
    chart_data.add_series('Avg Gross Hourly Wage (€)', (22.50, 18.50, 14.20))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.plots[0].has_data_labels = True
    add_axis_titles(chart, x_title='Employment Group', y_title='Hourly Wage (€)')

    colors = [BLUE, PURPLE, ORANGE]
    try:
        for idx, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(colors[idx])
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 5: Early Tracking & Educational Immobility
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "4. Early Tracking & Educational Immobility"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(PURPLE)
    
    p = tf.add_paragraph()
    p.text = "The Austrian education system separates children into different school tracks at the remarkably early age of 10. This practice strongly correlates educational outcomes with parental socioeconomic status, directly contradicting meritocratic ideals. A child of university graduates is roughly 4X more likely to obtain a degree than a child of parents with mandatory schooling."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(PURPLE)
    
    p = tf.add_paragraph()
    p.text = "• Delay binding student tracking to age 14 or 15 via comprehensive schools (Gesamtschule).\n• Implement a strict social-index-based funding model to allocate maximum resources to disadvantaged districts."
    p.font.size = Pt(14)

    # Chart: Column
    chart_data = CategoryChartData()
    chart_data.categories = ['Parents: Mandatory', 'Parents: Apprenticeship', 'Parents: Uni Degree']
    chart_data.add_series('Prob. Achieved Tertiary Degree (%)', (15, 28, 62))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.plots[0].has_data_labels = True
    add_axis_titles(chart, x_title='Parental Education', y_title='Child Probability (%)')
    
    try:
        series = chart.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = hex_to_rgb(PURPLE)
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 6: Spatial Divides (Scatter Plot)
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "5. Spatial Divides: Speculation vs. Urbanity"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BLUE)
    
    p = tf.add_paragraph()
    p.text = "Austria faces a severe geographic inequality challenge. Western Alpine regions (e.g., Vorarlberg, Salzburg) face extreme housing unaffordability due to tourism and speculation, diluting their high median incomes. Conversely, Vienna manages housing via public intervention and demonstrates economic resilience via a robust IT and tertiary sector (WIFO, 2024)."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BLUE)
    
    p = tf.add_paragraph()
    p.text = "• Institute strict vacancy taxes to halt real estate speculation in Western provinces.\n• Export Vienna's successful 'Gemeindebau' social housing model federally.\n• Relocate federal agencies outside the capital to stimulate rural local economies."
    p.font.size = Pt(14)

    # Chart: XY Scatter
    chart_data = XyChartData()
    series = chart_data.add_series('Regions')
    series.add_data_point(38000, 15.5) # Vorarlberg
    series.add_data_point(36500, 14.8) # Salzburg
    series.add_data_point(35000, 11.2) # Lower Austria
    series.add_data_point(31000, 10.5) # Vienna

    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data).chart
    add_axis_titles(chart, x_title='Median Annual Net Income (€)', y_title='Avg Rent Cost per sqm (€)')
    
    try:
        # Increase marker size and set color
        for idx, point in enumerate(chart.series[0].points):
            point.marker.size = 15
    except Exception:
        pass


    # -------------------------------------------------------------------------
    # SLIDE 7: The Regressive Tax of Inflation
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "6. The Regressive Tax of Inflation"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(RED)
    
    p = tf.add_paragraph()
    p.text = "Inflationary shocks are fundamentally regressive. The bottom 20% of earners spend over 40% of their income on non-discretionary necessities (housing and food). Consequently, their 'effective inflation rate' during price shocks is drastically higher than the top 20%, who spend a fraction of their income on basic needs."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(RED)
    
    p = tf.add_paragraph()
    p.text = "• Abandon broad-brush subsidies for targeted cash transfers.\n• Institute strict, temporary price caps on basic utility tariffs.\n• Decouple rent indexation from raw consumer price indices (CPI), tying it to median wage growth."
    p.font.size = Pt(14)

    # Chart: Line
    chart_data = CategoryChartData()
    chart_data.categories = ['21 Q4', '22 Q2', '22 Q4', '23 Q2', '23 Q4']
    chart_data.add_series('Bottom 20% Income', (4.5, 9.2, 11.8, 10.5, 6.8))
    chart_data.add_series('Top 20% Income', (4.0, 7.5, 9.1, 8.2, 5.5))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_axis_titles(chart, x_title='Quarter', y_title='Effective Inflation Rate (%)')

    try:
        chart.series[0].format.line.color.rgb = hex_to_rgb(RED)
        chart.series[1].format.line.color.rgb = hex_to_rgb(BLUE)
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 8: Political Inequality & "Diploma Democracy"
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "7. Political Inequality & 'Diploma Democracy'"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Findings / Key Argument:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BLUE)
    
    p = tf.add_paragraph()
    p.text = "Economic inequality translates relentlessly into political inequality. The Austrian National Council has evolved into a 'Diploma Democracy', heavily dominated by academics and affluent professionals. Marginalized and working-class groups abstain from voting due to alienation, ensuring their material needs are routinely ignored."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nPolicy Interventions:"
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BLUE)
    
    p = tf.add_paragraph()
    p.text = "• Explore reintroducing compulsory voting to perfectly mirror the socio-economic diversity of the nation.\n• Execute robust civic education pipelines in vocational schools.\n• Implement binding internal diversity quotas within political party recruitment for working-class backgrounds."
    p.font.size = Pt(14)

    # Chart: Clustered Column
    chart_data = CategoryChartData()
    chart_data.categories = ['Working Class', 'Academics / Professionals']
    chart_data.add_series('Share of Gen. Workforce', (35, 20))
    chart_data.add_series('Share of Parliament', (3, 60))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.plots[0].has_data_labels = True
    add_axis_titles(chart, x_title='Socio-Economic Group', y_title='Share (%)')
    
    try:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = hex_to_rgb(ORANGE)
        chart.series[1].format.fill.solid()
        chart.series[1].format.fill.fore_color.rgb = hex_to_rgb(PURPLE)
    except Exception:
        pass

    # -------------------------------------------------------------------------
    # SLIDE 9: Sources & References
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Sources & Academic References"
    
    # Adjust title size
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.clear()
    
    sources = [
        "[1] Oesterreichische Nationalbank (OeNB) (2024). Eurosystem Household Finance and Consumption Survey (HFCS) 2023.",
        "[2] PubMed (2022). A Tale of Integration? The Migrant Wealth Gap in Austria.",
        "[3] OECD Economic Surveys: Austria 2024.",
        "[4] Institute for Fiscal Studies (2023). Persistent low inequality despite compositional shifts in Austria.",
        "[5] Promoting social mobility in Austria - OECD.",
        "[6] Education and Training Monitor 2025 – Austria.",
        "[7] WIFO (2024). Second Year of Recession in Austria: Economic Development and Regional Disparities.",
        "[8] Statistics Austria. Economic output of most federal provinces almost unchanged in 2024.",
        "[9] Kontrast. Das sind die neuen Abgeordneten im Nationalrat.",
        "[10] ResearchGate (2023). Right-wing populism against diploma democracy: parliamentary elites in Austria."
    ]
    
    for src in sources:
        p = tf.add_paragraph()
        p.text = src
        p.font.size = Pt(13)
        p.space_after = Pt(4)

    # Save presentation
    file_name = "Austria_Inequality_Trends.pptx"
    prs.save(file_name)
    print(f"Presentation generated successfully: {file_name}")

if __name__ == "__main__":
    create_inequality_presentation()